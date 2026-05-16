import os
from spire.pdf import PdfDocument, PdfImageType
from pptx import Presentation
from pptx.util import Inches

pdf_path = r"c:\Code\SentimentFlow\slidev\slides-export.pdf"
pptx_path = r"c:\Code\SentimentFlow\slidev\slides-export.pptx"

print(f"Converting {pdf_path} to PPTX using Spire.PDF...")

try:
    doc = PdfDocument()
    doc.LoadFromFile(pdf_path)

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(doc.Pages.Count):
        img_temp_path = f"temp_slide_{i}.png"
        with doc.SaveAsImage(i, PdfImageType.Bitmap) as image:
            image.Save(img_temp_path)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slide.shapes.add_picture(img_temp_path, 0, 0, width=slide_width, height=slide_height)
        
        os.remove(img_temp_path)

    doc.Close()
    prs.save(pptx_path)
    print(f"Successfully saved PPTX to {pptx_path}")
except Exception as e:
    print(f"Error during conversion: {e}")
    import sys
    sys.exit(1)
