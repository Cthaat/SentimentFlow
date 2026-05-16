import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

pdf_path = r'c:\Code\SentimentFlow\slidev\slides-export.pdf'
pptx_path = r'c:\Code\SentimentFlow\slidev\slides-export.pptx'

print(f"Converting {pdf_path} to PPTX...")

try:
    # Convert PDF to list of images
    # Note: poppler must be installed for pdf2image to work. 
    # If it fails, I might need to check for poppler.
    images = convert_from_path(pdf_path)

    prs = Presentation()
    
    # Remove default slide layout
    # prs.slide_layouts[6] is usually a blank slide
    blank_slide_layout = prs.slide_layouts[6]

    for i, image in enumerate(images):
        # Save image temporarily
        img_temp_path = f'temp_slide_{i}.png'
        image.save(img_temp_path, 'PNG')
        
        # Add slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add image to slide
        # Calculate dimensions to fill the slide (roughly)
        # Default slide size is 10 x 7.5 inches
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        slide.shapes.add_picture(img_temp_path, 0, 0, width=slide_width, height=slide_height)
        
        # Remove temp image
        os.remove(img_temp_path)

    prs.save(pptx_path)
    print(f"Successfully saved PPTX to {pptx_path}")
except Exception as e:
    print(f"Error during conversion: {e}")
    import sys
    sys.exit(1)
